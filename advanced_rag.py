"""
Advanced RAG (Retrieval Augmented Generation) system for Ollama Workbench.

This module provides enhanced document processing, storage, and retrieval capabilities
for more effective RAG implementations. It includes:

1. Advanced document processing with semantic chunking
2. Hybrid search combining semantic and keyword search
3. Document categorization and metadata management
4. Customizable retrieval strategies
5. Integration with multiple vector databases
6. Query expansion and rewriting
7. Result re-ranking
"""

import os
import json
import shutil
import time
import uuid
import re
import logging
import numpy as np
import pandas as pd
from datetime import datetime
from typing import List, Dict, Any, Optional, Union, Tuple, Callable
from sklearn.metrics.pairwise import cosine_similarity
from collections import defaultdict

# Document processing
import PyPDF2
from bs4 import BeautifulSoup
import markdown
import docx2txt
import pptx

# Vector databases
import chromadb
from chromadb.config import Settings

# Embedding generation
import ollama
from ollama_utils import get_ollama_client, call_ollama_endpoint
import openai

logger = logging.getLogger('advanced_rag')

# Constants
DOCUMENT_STORE_DIR = os.path.join("data", "document_store")
EMBEDDING_CACHE_DIR = os.path.join("data", "embedding_cache")
CONFIG_DIR = os.path.join("data", "rag_config")

# Ensure directories exist
os.makedirs(DOCUMENT_STORE_DIR, exist_ok=True)
os.makedirs(EMBEDDING_CACHE_DIR, exist_ok=True) 
os.makedirs(CONFIG_DIR, exist_ok=True)

# Default configuration
DEFAULT_CONFIG = {
    "chunking": {
        "method": "semantic",  # semantic, fixed, sliding, paragraph
        "size": 1000,          # characters or tokens depending on config
        "overlap": 200,        # overlap between chunks
        "unit": "characters",  # characters or tokens
        "respect_boundaries": True,  # respect sentence/paragraph boundaries
    },
    "embedding": {
        "provider": "ollama",  # ollama, openai, custom
        "model": "nomic-embed-text",  # default model for embeddings
        "batch_size": 10,      # number of documents to embed at once
        "embedding_dim": 768,  # dimension of embeddings
        "cache_embeddings": True,
    },
    "retrieval": {
        "search_type": "hybrid",  # semantic, keyword, hybrid
        "semantic_weight": 0.7,   # weight for semantic search in hybrid (0-1)
        "keyword_weight": 0.3,    # weight for keyword search in hybrid (0-1)
        "top_k": 5,              # number of documents to retrieve
        "min_similarity": 0.2,    # minimum similarity score to include document
        "use_reranking": True,    # whether to rerank results
        "reranker": "default",    # reranking method to use
    },
    "vector_db": {
        "backend": "chroma",   # chroma, faiss, pgvector, milvus
        "persist_directory": os.path.join(DOCUMENT_STORE_DIR, "vectordb"),
        "collection_name": "documents",
    },
    "query_processing": {
        "expand_queries": True,  # whether to expand queries
        "expansion_method": "llm",  # llm, thesaurus, none
        "rewrite_queries": True,  # whether to rewrite queries
        "rewrite_method": "llm",  # llm, template, none
    },
    "document_processing": {
        "extract_metadata": True,  # extract metadata from documents
        "categorize_documents": True,  # categorize documents
        "summarize_documents": True,  # generate summaries of documents
        "detect_language": True,  # detect document language
    }
}

class DocumentTypes:
    """Enumeration of supported document types."""
    TEXT = "text"
    PDF = "pdf"
    HTML = "html"
    MARKDOWN = "markdown"
    DOCX = "docx"
    PPTX = "pptx"
    JSON = "json"
    CSV = "csv"
    CODE = "code"
    UNKNOWN = "unknown"
    
    @classmethod
    def detect_type(cls, file_path: str) -> str:
        """Detect document type from file extension."""
        if not file_path:
            return cls.UNKNOWN
            
        extension = file_path.lower().split('.')[-1] if '.' in file_path else ''
        
        if extension in ['txt', 'text']:
            return cls.TEXT
        elif extension == 'pdf':
            return cls.PDF
        elif extension in ['html', 'htm']:
            return cls.HTML
        elif extension in ['md', 'markdown']:
            return cls.MARKDOWN
        elif extension == 'docx':
            return cls.DOCX
        elif extension == 'pptx':
            return cls.PPTX
        elif extension == 'json':
            return cls.JSON
        elif extension == 'csv':
            return cls.CSV
        elif extension in ['py', 'js', 'java', 'cpp', 'c', 'cs', 'go', 'rs', 'php', 'rb', 'swift', 'kt']:
            return cls.CODE
        else:
            return cls.UNKNOWN


class Document:
    """Class representing a document with its metadata and content."""
    
    def __init__(
        self, 
        content: str,
        doc_id: Optional[str] = None,
        title: Optional[str] = None,
        source: Optional[str] = None,
        doc_type: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        chunks: Optional[List[Dict[str, Any]]] = None,
        embedding: Optional[np.ndarray] = None,
        summary: Optional[str] = None,
        category: Optional[str] = None,
    ):
        self.doc_id = doc_id or str(uuid.uuid4())
        self.content = content
        self.title = title or f"Document-{self.doc_id[:8]}"
        self.source = source
        self.doc_type = doc_type or DocumentTypes.UNKNOWN
        self.metadata = metadata or {}
        self.chunks = chunks or []
        self.embedding = embedding
        self.summary = summary
        self.category = category
        self.created_at = datetime.now().isoformat()
        self.updated_at = self.created_at
        
    def to_dict(self) -> Dict[str, Any]:
        """Convert document to dictionary for serialization."""
        return {
            "doc_id": self.doc_id,
            "title": self.title,
            "content": self.content,
            "source": self.source,
            "doc_type": self.doc_type,
            "metadata": self.metadata,
            "chunks": self.chunks,
            "embedding": self.embedding.tolist() if isinstance(self.embedding, np.ndarray) else self.embedding,
            "summary": self.summary,
            "category": self.category,
            "created_at": self.created_at,
            "updated_at": self.updated_at
        }
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'Document':
        """Create document from dictionary."""
        embedding = data.get('embedding')
        if embedding is not None and isinstance(embedding, list):
            embedding = np.array(embedding)
            
        return cls(
            content=data.get('content', ''),
            doc_id=data.get('doc_id'),
            title=data.get('title'),
            source=data.get('source'),
            doc_type=data.get('doc_type'),
            metadata=data.get('metadata'),
            chunks=data.get('chunks'),
            embedding=embedding,
            summary=data.get('summary'),
            category=data.get('category')
        )
    
    def update_metadata(self, metadata: Dict[str, Any]) -> None:
        """Update document metadata."""
        self.metadata.update(metadata)
        self.updated_at = datetime.now().isoformat()


class DocumentChunker:
    """Class for chunking documents using various strategies."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize the document chunker with the given configuration."""
        self.config = config or DEFAULT_CONFIG["chunking"]
        self.semantic_chunking_enabled = True  # Flag to control semantic chunking
        
    def chunk_document(self, document: Document) -> List[Dict[str, Any]]:
        """Chunk a document based on the configured strategy."""
        method = self.config.get("method", "fixed")
        
        if method == "semantic" and self.semantic_chunking_enabled:
            try:
                chunks = self._semantic_chunking(document)
                logger.info(f"Semantic chunking created {len(chunks)} chunks")
            except Exception as e:
                logger.error(f"Semantic chunking failed: {str(e)}. Falling back to paragraph chunking.")
                chunks = self._paragraph_chunking(document)
        elif method == "sliding":
            chunks = self._sliding_window_chunking(document)
        elif method == "paragraph":
            chunks = self._paragraph_chunking(document)
        else:  # Default to fixed-size chunking
            chunks = self._fixed_size_chunking(document)
            
        # Add metadata to chunks
        for i, chunk in enumerate(chunks):
            chunk["doc_id"] = document.doc_id
            chunk["chunk_id"] = f"{document.doc_id}-{i}"
            chunk["document_title"] = document.title
            chunk["position"] = i
            
        return chunks
    
    def _fixed_size_chunking(self, document: Document) -> List[Dict[str, Any]]:
        """Split text into fixed-size chunks."""
        text = document.content
        size = self.config.get("size", 1000)
        overlap = self.config.get("overlap", 200)
        respect_boundaries = self.config.get("respect_boundaries", True)
        
        chunks = []
        
        if respect_boundaries:
            # Split by sentences first
            sentences = re.split(r'(?<=[.!?])\s+', text)
            current_chunk = ""
            
            for sentence in sentences:
                if len(current_chunk) + len(sentence) > size and current_chunk:
                    chunks.append({"content": current_chunk.strip()})
                    # Keep overlap if possible
                    words = current_chunk.split()
                    overlap_word_count = min(len(words), int(overlap / 5))  # Approximate words in overlap
                    if overlap_word_count > 0:
                        current_chunk = " ".join(words[-overlap_word_count:]) + " "
                    else:
                        current_chunk = ""
                
                current_chunk += sentence + " "
            
            # Add the last chunk if it's not empty
            if current_chunk.strip():
                chunks.append({"content": current_chunk.strip()})
        else:
            # Simple character-based chunking
            for i in range(0, len(text), size - overlap):
                chunk_text = text[i:i + size]
                chunks.append({"content": chunk_text.strip()})
                
        return chunks
    
    def _sliding_window_chunking(self, document: Document) -> List[Dict[str, Any]]:
        """Split text using a sliding window approach with more granular control."""
        text = document.content
        size = self.config.get("size", 1000)
        overlap = self.config.get("overlap", 200)
        step = size - overlap
        
        chunks = []
        
        for i in range(0, len(text), step):
            chunk_text = text[i:i + size]
            if len(chunk_text.strip()) > 0:
                chunks.append({"content": chunk_text.strip()})
                
        return chunks
    
    def _paragraph_chunking(self, document: Document) -> List[Dict[str, Any]]:
        """Split text by paragraphs and then combine into chunks of appropriate size."""
        text = document.content
        size = self.config.get("size", 1000)
        
        # Split by paragraphs (double newlines)
        paragraphs = re.split(r'\n\s*\n', text)
        
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            if len(current_chunk) + len(paragraph) > size and current_chunk:
                chunks.append({"content": current_chunk.strip()})
                current_chunk = ""
            
            current_chunk += paragraph + "\n\n"
        
        # Add the last chunk if it's not empty
        if current_chunk.strip():
            chunks.append({"content": current_chunk.strip()})
            
        return chunks
    
    def _semantic_chunking(self, document: Document) -> List[Dict[str, Any]]:
        """
        Split text into semantically coherent chunks.
        
        This uses a more advanced approach that tries to keep related content together by:
        1. First splitting into paragraphs
        2. Then calculating semantic similarity between paragraphs
        3. Combining similar paragraphs into chunks up to the size limit
        """
        from ollama_utils import get_ollama_client, call_ollama_endpoint
        import numpy as np
        from sklearn.metrics.pairwise import cosine_similarity
        
        # Set the max size for chunks (in characters)
        max_chunk_size = self.config.get("size", 1000)
        text = document.content
        
        # Step 1: Split the text into paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]
        
        # If very few paragraphs, just use paragraph chunking
        if len(paragraphs) <= 2:
            return self._paragraph_chunking(document)
        
        # Step 2: Calculate embeddings for each paragraph
        try:
            # Initialize embedder
            model = "nomic-embed-text"  # Default embedding model
            client = get_ollama_client()
            
            # Get embeddings for each paragraph
            paragraph_embeddings = []
            for paragraph in paragraphs:
                try:
                    if client:
                        response = client.embeddings(model=model, prompt=paragraph)
                        embedding = np.array(response['embedding'])
                    else:
                        # Fallback to direct API call
                        import requests
                        response = requests.post(
                            "http://localhost:11434/api/embeddings",
                            json={"model": model, "prompt": paragraph}
                        )
                        response.raise_for_status()
                        embedding = np.array(response.json()['embedding'])
                    paragraph_embeddings.append(embedding)
                except Exception as e:
                    logger.warning(f"Error getting embedding for paragraph: {str(e)}")
                    # If embedding fails, use an array of zeros as a placeholder
                    paragraph_embeddings.append(np.zeros(768))  # Typical embedding size
            
            # Step 3: Calculate similarity matrix between paragraphs
            similarity_matrix = cosine_similarity(paragraph_embeddings)
            
            # Step 4: Form chunks by combining similar paragraphs
            chunks = []
            current_chunk_paragraphs = [0]  # Start with the first paragraph
            current_chunk_text = paragraphs[0]
            visited = {0}  # Keep track of paragraphs already in chunks
            
            # Process each paragraph (starting from the second one)
            for i in range(1, len(paragraphs)):
                if i in visited:
                    continue
                    
                # Find the paragraph in the current chunk with highest similarity to paragraph i
                similarities = [similarity_matrix[j][i] for j in current_chunk_paragraphs]
                max_similarity = max(similarities) if similarities else 0
                
                # If the similarity is high enough and adding won't exceed size limit, add to current chunk
                new_chunk_text = current_chunk_text + "\n\n" + paragraphs[i]
                if max_similarity > 0.5 and len(new_chunk_text) <= max_chunk_size:
                    current_chunk_paragraphs.append(i)
                    current_chunk_text = new_chunk_text
                    visited.add(i)
                else:
                    # Save current chunk and start a new one
                    chunks.append({"content": current_chunk_text.strip()})
                    current_chunk_paragraphs = [i]
                    current_chunk_text = paragraphs[i]
                    visited.add(i)
            
            # Add the last chunk if it's not empty
            if current_chunk_text.strip():
                chunks.append({"content": current_chunk_text.strip()})
                
            # Process remaining paragraphs (if any weren't visited yet)
            remaining = [i for i in range(len(paragraphs)) if i not in visited]
            for i in remaining:
                # Try to find the most similar existing chunk
                max_similarity = 0
                best_chunk = -1
                
                for j, chunk in enumerate(chunks):
                    # Find paragraphs in this chunk
                    chunk_text = chunk["content"]
                    for k in range(len(paragraphs)):
                        if k in visited and paragraphs[k] in chunk_text:
                            # Calculate similarity between paragraph i and paragraph k
                            similarity = similarity_matrix[i][k]
                            if similarity > max_similarity:
                                max_similarity = similarity
                                best_chunk = j
                
                # If found a similar chunk and adding won't exceed size limit, add to that chunk
                if best_chunk >= 0 and max_similarity > 0.5:
                    new_content = chunks[best_chunk]["content"] + "\n\n" + paragraphs[i]
                    if len(new_content) <= max_chunk_size:
                        chunks[best_chunk]["content"] = new_content
                        visited.add(i)
                        continue
                
                # Otherwise, create a new chunk
                chunks.append({"content": paragraphs[i].strip()})
                visited.add(i)
                
            # Ensure all paragraphs are included
            if len(visited) < len(paragraphs):
                # Fall back to paragraph chunking if our algorithm missed any paragraphs
                logger.warning(f"Semantic chunking missed some paragraphs. Falling back to paragraph chunking.")
                return self._paragraph_chunking(document)
                
            return chunks
            
        except Exception as e:
            logger.error(f"Error in semantic chunking: {str(e)}")
            # Fall back to paragraph chunking if any errors
            return self._paragraph_chunking(document)


class DocumentProcessor:
    """Advanced document processor for handling various document types."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize the document processor with the given configuration."""
        self.config = config or DEFAULT_CONFIG["document_processing"]
        self.chunker = DocumentChunker(config["chunking"] if config else None)
        
    def process_file(self, file_path: str) -> Document:
        """Process a file and return a Document object."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        doc_type = DocumentTypes.detect_type(file_path)
        content = self._extract_text_from_file(file_path, doc_type)
        
        title = os.path.basename(file_path)
        
        doc = Document(
            content=content,
            title=title,
            source=file_path,
            doc_type=doc_type
        )
        
        # Process document further based on configuration
        if self.config.get("extract_metadata", True):
            doc.metadata = self._extract_metadata(doc)
            
        if self.config.get("categorize_documents", True):
            doc.category = self._categorize_document(doc)
            
        if self.config.get("summarize_documents", True):
            doc.summary = self._summarize_document(doc)
            
        if self.config.get("detect_language", True):
            doc.metadata["language"] = self._detect_language(doc)
        
        # Chunk the document
        doc.chunks = self.chunker.chunk_document(doc)
        
        return doc
    
    def process_text(self, text: str, title: str = None, source: str = None, doc_type: str = DocumentTypes.TEXT) -> Document:
        """Process text and return a Document object."""
        doc = Document(
            content=text,
            title=title or "Text Document",
            source=source,
            doc_type=doc_type
        )
        
        # Process document further based on configuration
        if self.config.get("extract_metadata", True):
            doc.metadata = self._extract_metadata(doc)
            
        if self.config.get("categorize_documents", True):
            doc.category = self._categorize_document(doc)
            
        if self.config.get("summarize_documents", True):
            doc.summary = self._summarize_document(doc)
            
        if self.config.get("detect_language", True):
            doc.metadata["language"] = self._detect_language(doc)
        
        # Chunk the document
        doc.chunks = self.chunker.chunk_document(doc)
        
        return doc
    
    def _extract_text_from_file(self, file_path: str, doc_type: str) -> str:
        """Extract text from a file based on its type."""
        try:
            if doc_type == DocumentTypes.TEXT:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
                    
            elif doc_type == DocumentTypes.PDF:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n\n"
                    return text
                    
            elif doc_type == DocumentTypes.HTML:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.extract()
                    return soup.get_text(separator='\n')
                    
            elif doc_type == DocumentTypes.MARKDOWN:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    md_text = f.read()
                    html = markdown.markdown(md_text)
                    soup = BeautifulSoup(html, 'html.parser')
                    return soup.get_text(separator='\n')
                    
            elif doc_type == DocumentTypes.DOCX:
                return docx2txt.process(file_path)
                
            elif doc_type == DocumentTypes.PPTX:
                presentation = pptx.Presentation(file_path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n\n"
                return text
                
            elif doc_type == DocumentTypes.JSON:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2)
                    
            elif doc_type == DocumentTypes.CSV:
                df = pd.read_csv(file_path)
                return df.to_string()
                
            elif doc_type == DocumentTypes.CODE:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
                    
            else:
                # Try to read as text for unknown types
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
                    
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return f"Error extracting text: {str(e)}"
    
    def _extract_metadata(self, document: Document) -> Dict[str, Any]:
        """Extract metadata from a document."""
        metadata = {}
        
        # Basic metadata
        metadata["doc_type"] = document.doc_type
        metadata["file_size"] = len(document.content)
        metadata["word_count"] = len(document.content.split())
        metadata["created_at"] = document.created_at
        
        # Type-specific metadata
        if document.doc_type == DocumentTypes.PDF:
            # In a more complete implementation, extract PDF metadata
            pass
        elif document.doc_type == DocumentTypes.HTML:
            # Extract HTML metadata
            pass
        
        return metadata
    
    def _categorize_document(self, document: Document) -> str:
        """Categorize a document based on its content."""
        # This would normally use an LLM or classifier
        # For now, use a simple heuristic based on document type
        return document.doc_type
    
    def _summarize_document(self, document: Document) -> str:
        """Generate a summary of the document."""
        # This would normally use an LLM for summarization
        # For now, return the first 200 characters as a placeholder
        text = document.content.strip()
        if len(text) <= 200:
            return text
        return text[:200] + "..."
    
    def _detect_language(self, document: Document) -> str:
        """Detect the language of the document."""
        # This would normally use a language detection library
        # For now, assume English
        return "en"


class EmbeddingGenerator:
    """Class for generating embeddings from text using various providers."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize the embedding generator with the given configuration."""
        self.config = config or DEFAULT_CONFIG["embedding"]
        self.cache = {}
        self.cache_file = os.path.join(EMBEDDING_CACHE_DIR, "embedding_cache.json")
        
        # Load embedding cache if enabled
        if self.config.get("cache_embeddings", True):
            self._load_cache()
    
    def _load_cache(self) -> None:
        """Load embedding cache from file."""
        if os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r') as f:
                    cache_data = json.load(f)
                    
                # Convert lists back to numpy arrays
                for key, value in cache_data.items():
                    self.cache[key] = np.array(value)
                    
                logger.info(f"Loaded {len(self.cache)} embeddings from cache")
            except Exception as e:
                logger.error(f"Error loading embedding cache: {str(e)}")
                self.cache = {}
    
    def _save_cache(self) -> None:
        """Save embedding cache to file."""
        if not self.config.get("cache_embeddings", True):
            return
            
        try:
            cache_data = {}
            
            # Convert numpy arrays to lists for JSON serialization
            for key, value in self.cache.items():
                if isinstance(value, np.ndarray):
                    cache_data[key] = value.tolist()
                else:
                    cache_data[key] = value
                    
            with open(self.cache_file, 'w') as f:
                json.dump(cache_data, f)
                
            logger.info(f"Saved {len(self.cache)} embeddings to cache")
        except Exception as e:
            logger.error(f"Error saving embedding cache: {str(e)}")
    
    def get_embedding(self, text: str) -> np.ndarray:
        """Get embedding for a piece of text."""
        # Check cache first
        if self.config.get("cache_embeddings", True):
            # Use a truncated version of the text as a key
            cache_key = text[:1000]
            if cache_key in self.cache:
                return self.cache[cache_key]
        
        # Generate embedding based on provider
        provider = self.config.get("provider", "ollama")
        
        if provider == "ollama":
            embedding = self._get_ollama_embedding(text)
        elif provider == "openai":
            embedding = self._get_openai_embedding(text)
        else:
            # Default to Ollama
            embedding = self._get_ollama_embedding(text)
        
        # Cache the result
        if self.config.get("cache_embeddings", True) and embedding is not None:
            self.cache[cache_key] = embedding
            # Periodically save cache
            if len(self.cache) % 10 == 0:
                self._save_cache()
        
        return embedding
    
    def get_embeddings(self, texts: List[str]) -> List[np.ndarray]:
        """Get embeddings for multiple pieces of text."""
        # Process in batches
        batch_size = self.config.get("batch_size", 10)
        embeddings = []
        
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            batch_embeddings = [self.get_embedding(text) for text in batch]
            embeddings.extend(batch_embeddings)
        
        return embeddings
    
    def embed_document(self, document: Document) -> Document:
        """Generate embeddings for a document and its chunks."""
        # Generate document-level embedding
        document.embedding = self.get_embedding(document.content)
        
        # Generate embeddings for each chunk
        for chunk in document.chunks:
            chunk["embedding"] = self.get_embedding(chunk["content"])
        
        return document
    
    def _get_ollama_embedding(self, text: str) -> np.ndarray:
        """Get embedding from Ollama API."""
        model = self.config.get("model", "nomic-embed-text")
        
        try:
            # Try using the Ollama client class first
            client = get_ollama_client()
            if client:
                response = client.embeddings(model=model, prompt=text)
                return np.array(response['embedding'])
            
            # Fall back to direct API call
            url = "http://localhost:11434/api/embeddings"
            payload = {
                "model": model,
                "prompt": text
            }
            
            headers = {
                "Content-Type": "application/json"
            }
            
            response = requests.post(url, json=payload, headers=headers)
            if response.status_code == 200:
                return np.array(response.json()['embedding'])
            else:
                logger.error(f"Error getting Ollama embedding: {response.text}")
                return None
        except Exception as e:
            logger.error(f"Error getting Ollama embedding: {str(e)}")
            return None
    
    def _get_openai_embedding(self, text: str) -> np.ndarray:
        """Get embedding from OpenAI API."""
        model = self.config.get("model", "text-embedding-3-small")
        
        try:
            response = openai.Embedding.create(
                input=text,
                model=model
            )
            
            embedding = response['data'][0]['embedding']
            return np.array(embedding)
        except Exception as e:
            logger.error(f"Error getting OpenAI embedding: {str(e)}")
            return None


class QueryProcessor:
    """Class for processing and transforming queries."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize the query processor with the given configuration."""
        self.config = config or DEFAULT_CONFIG["query_processing"]
    
    def process_query(self, query: str) -> List[str]:
        """Process and expand a query."""
        processed_queries = [query]
        
        # Expand query if configured
        if self.config.get("expand_queries", True):
            expansion_method = self.config.get("expansion_method", "llm")
            if expansion_method == "llm":
                expanded_queries = self._expand_query_llm(query)
                processed_queries.extend(expanded_queries)
            elif expansion_method == "thesaurus":
                expanded_queries = self._expand_query_thesaurus(query)
                processed_queries.extend(expanded_queries)
        
        # Rewrite query if configured
        if self.config.get("rewrite_queries", True):
            rewrite_method = self.config.get("rewrite_method", "llm")
            if rewrite_method == "llm":
                rewritten_query = self._rewrite_query_llm(query)
                if rewritten_query and rewritten_query != query:
                    processed_queries.append(rewritten_query)
            elif rewrite_method == "template":
                rewritten_query = self._rewrite_query_template(query)
                if rewritten_query and rewritten_query != query:
                    processed_queries.append(rewritten_query)
        
        return processed_queries
    
    def _expand_query_llm(self, query: str) -> List[str]:
        """Expand a query using an LLM to generate related queries."""
        try:
            # Prepare prompt for query expansion
            prompt = f"""Generate 3 alternative search queries related to the following query:
"{query}"

The alternative queries should:
1. Capture the same information need from different angles
2. Use synonyms and related concepts
3. Vary in specificity (more general and more specific variations)

Return ONLY the queries, one per line, with no numbering or other text.
"""
            
            # Use ollama for query expansion
            from ollama_utils import call_ollama_endpoint
            
            # Default to a fast, efficient model
            model = "llama3"  # Could be configured in settings
            
            # Call the model
            response, _, _, _ = call_ollama_endpoint(
                model=model,
                prompt=prompt,
                temperature=0.7,
                max_tokens=200
            )
            
            # Parse response into separate queries
            expanded_queries = []
            for line in response.strip().split('\n'):
                line = line.strip()
                if line and line != query and not line.startswith('#') and not line.startswith('-'):
                    # Clean up any quotes or artifacts
                    line = line.strip('"\'')
                    expanded_queries.append(line)
            
            # If we didn't get any valid expansions, fall back to rule-based
            if not expanded_queries:
                if not query.lower().startswith("what is"):
                    expanded_queries.append(f"What is {query}")
                if not query.lower().startswith("how to"):
                    expanded_queries.append(f"How to {query}")
            
            return expanded_queries
        except Exception as e:
            logger.error(f"Error in query expansion with LLM: {str(e)}")
            # Fallback to rule-based expansion
            expanded_queries = []
            if not query.lower().startswith("what is"):
                expanded_queries.append(f"What is {query}")
            if not query.lower().startswith("how to"):
                expanded_queries.append(f"How to {query}")
            return expanded_queries
    
    def _expand_query_thesaurus(self, query: str) -> List[str]:
        """Expand a query using synonyms for key terms."""
        # Simple synonym sets for common terms
        synonyms = {
            "create": ["make", "build", "develop", "generate", "produce"],
            "fix": ["repair", "resolve", "mend", "correct", "address"],
            "error": ["bug", "issue", "problem", "defect", "fault"],
            "guide": ["tutorial", "instructions", "walkthrough", "handbook", "manual"],
            "best": ["optimal", "ideal", "top", "excellent", "superior"],
            "fast": ["quick", "rapid", "swift", "speedy", "high-performance"],
            "example": ["sample", "instance", "illustration", "demonstration", "case"],
            "important": ["essential", "crucial", "critical", "vital", "significant"],
            "feature": ["functionality", "capability", "aspect", "component", "attribute"]
        }
        
        # Tokenize the query into words
        words = query.lower().split()
        expanded_queries = []
        
        # Generate expanded queries by replacing one word at a time
        for i, word in enumerate(words):
            if word in synonyms:
                for synonym in synonyms[word]:
                    new_words = words.copy()
                    new_words[i] = synonym
                    expanded_query = " ".join(new_words)
                    if expanded_query != query.lower():
                        expanded_queries.append(expanded_query)
        
        return expanded_queries[:3]  # Limit to 3 expansions
    
    def _rewrite_query_llm(self, query: str) -> str:
        """Rewrite a query using an LLM to make it more effective for retrieval."""
        try:
            # Prepare prompt for query rewriting
            prompt = f"""Rewrite the following search query to make it more effective for retrieving relevant information:
"{query}"

A good search query should:
1. Be clear and specific
2. Include important keywords
3. Be focused on retrieving factual information
4. Avoid unnecessary words

Return ONLY the rewritten query with no additional text or explanation.
"""
            
            # Use ollama for query rewriting
            from ollama_utils import call_ollama_endpoint
            
            # Default to a fast, efficient model
            model = "llama3"  # Could be configured in settings
            
            # Call the model
            response, _, _, _ = call_ollama_endpoint(
                model=model,
                prompt=prompt,
                temperature=0.3,  # Lower temperature for more focused responses
                max_tokens=100
            )
            
            # Clean up response
            rewritten_query = response.strip().strip('"\'')
            
            # Only return if it's actually different
            if rewritten_query.lower() != query.lower():
                return rewritten_query
            return None
        except Exception as e:
            logger.error(f"Error in query rewriting with LLM: {str(e)}")
            return None
    
    def _rewrite_query_template(self, query: str) -> str:
        """Rewrite a query using templates."""
        # Simple template-based rewriting
        query = query.strip()
        
        # Remove question marks
        if query.endswith('?'):
            return query[:-1]
        
        return None


class DocumentStore:
    """Class for storing and retrieving documents."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize the document store with the given configuration."""
        self.config = config or {
            **DEFAULT_CONFIG["vector_db"],
            **DEFAULT_CONFIG["retrieval"]
        }
        
        self.embedder = EmbeddingGenerator(config["embedding"] if config else None)
        self.query_processor = QueryProcessor(config["query_processing"] if config else None)
        
        # Initialize vector database
        self.backend = self.config.get("backend", "chroma")
        self.collection_name = self.config.get("collection_name", "documents")
        self.persist_directory = self.config.get("persist_directory", os.path.join(DOCUMENT_STORE_DIR, "vectordb"))
        
        # Ensure directory exists
        os.makedirs(self.persist_directory, exist_ok=True)
        
        # Create persistent client
        if self.backend == "chroma":
            self.db_client = chromadb.PersistentClient(path=self.persist_directory)
            
            # Create or get collection
            try:
                self.collection = self.db_client.get_collection(name=self.collection_name)
                logger.info(f"Using existing collection: {self.collection_name}")
            except:
                self.collection = self.db_client.create_collection(name=self.collection_name)
                logger.info(f"Created new collection: {self.collection_name}")
        else:
            # For other backends, implement similar initialization
            logger.warning(f"Vector database backend '{self.backend}' not implemented, defaulting to file-based store")
            self.db_client = None
            self.collection = None
            
        # Document metadata store (for document-level operations)
        self.document_metadata_file = os.path.join(self.persist_directory, "document_metadata.json")
        self.document_metadata = self._load_document_metadata()
    
    def _load_document_metadata(self) -> Dict[str, Any]:
        """Load document metadata from file."""
        if os.path.exists(self.document_metadata_file):
            try:
                with open(self.document_metadata_file, 'r') as f:
                    return json.load(f)
            except Exception as e:
                logger.error(f"Error loading document metadata: {str(e)}")
                return {}
        else:
            return {}
    
    def _save_document_metadata(self) -> None:
        """Save document metadata to file."""
        try:
            with open(self.document_metadata_file, 'w') as f:
                json.dump(self.document_metadata, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving document metadata: {str(e)}")
    
    def add_document(self, document: Document) -> str:
        """Add a document to the store."""
        # Generate embeddings if not already present
        if document.embedding is None or not any([chunk.get("embedding") is not None for chunk in document.chunks]):
            document = self.embedder.embed_document(document)
        
        # Add to vector database if using one
        if self.collection:
            # Add each chunk as a separate item
            ids = []
            embeddings = []
            metadatas = []
            documents = []
            
            for chunk in document.chunks:
                chunk_id = chunk["chunk_id"]
                ids.append(chunk_id)
                
                # Convert numpy array to list for ChromaDB
                embedding = chunk["embedding"].tolist() if isinstance(chunk["embedding"], np.ndarray) else chunk["embedding"]
                embeddings.append(embedding)
                
                # Prepare metadata
                metadata = {
                    "doc_id": document.doc_id,
                    "chunk_id": chunk_id,
                    "document_title": document.title,
                    "position": chunk.get("position", 0),
                    "category": document.category,
                    "source": document.source,
                    "doc_type": document.doc_type
                }
                metadatas.append(metadata)
                
                # Document content
                documents.append(chunk["content"])
            
            # Add to collection
            self.collection.add(
                ids=ids,
                embeddings=embeddings,
                metadatas=metadatas,
                documents=documents
            )
        
        # Store document metadata
        self.document_metadata[document.doc_id] = {
            "doc_id": document.doc_id,
            "title": document.title,
            "source": document.source,
            "doc_type": document.doc_type,
            "summary": document.summary,
            "category": document.category,
            "metadata": document.metadata,
            "created_at": document.created_at,
            "updated_at": document.updated_at,
            "chunk_count": len(document.chunks)
        }
        
        # Save document metadata
        self._save_document_metadata()
        
        return document.doc_id
    
    def get_document(self, doc_id: str) -> Optional[Document]:
        """Retrieve a document by ID."""
        if doc_id not in self.document_metadata:
            return None
        
        # Get document metadata
        doc_metadata = self.document_metadata[doc_id]
        
        # Get document chunks from vector database
        if self.collection:
            results = self.collection.get(
                where={"doc_id": doc_id}
            )
            
            chunks = []
            for i in range(len(results['ids'])):
                chunk = {
                    "chunk_id": results['ids'][i],
                    "content": results['documents'][i],
                    "position": results['metadatas'][i].get("position", 0)
                }
                chunks.append(chunk)
            
            # Sort chunks by position
            chunks.sort(key=lambda x: x["position"])
            
            # Reconstruct full content
            content = "\n\n".join([chunk["content"] for chunk in chunks])
            
            # Create document
            document = Document(
                content=content,
                doc_id=doc_id,
                title=doc_metadata.get("title", ""),
                source=doc_metadata.get("source", ""),
                doc_type=doc_metadata.get("doc_type", DocumentTypes.UNKNOWN),
                metadata=doc_metadata.get("metadata", {}),
                chunks=chunks,
                summary=doc_metadata.get("summary", ""),
                category=doc_metadata.get("category", "")
            )
            
            return document
        else:
            # If not using a vector database, we can't reconstruct the document
            logger.warning(f"Cannot reconstruct document {doc_id} without vector database")
            return None
    
    def delete_document(self, doc_id: str) -> bool:
        """Delete a document from the store."""
        if doc_id not in self.document_metadata:
            return False
        
        # Delete from vector database
        if self.collection:
            self.collection.delete(
                where={"doc_id": doc_id}
            )
        
        # Delete from metadata
        del self.document_metadata[doc_id]
        
        # Save metadata
        self._save_document_metadata()
        
        return True
    
    def list_documents(self) -> List[Dict[str, Any]]:
        """List all documents in the store."""
        return list(self.document_metadata.values())
    
    def search(self, query: str, top_k: int = None, min_similarity: float = None, filter_criteria: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        Search for documents matching the query.
        
        Args:
            query (str): The search query.
            top_k (int, optional): The number of results to return. Defaults to config setting.
            min_similarity (float, optional): Minimum similarity score. Defaults to config setting.
            filter_criteria (Dict[str, Any], optional): Criteria to filter results.
            
        Returns:
            List[Dict[str, Any]]: List of search results with content and metadata.
        """
        # Process query
        processed_queries = self.query_processor.process_query(query)
        
        # Use configured or provided values
        top_k = top_k or self.config.get("top_k", 5)
        min_similarity = min_similarity or self.config.get("min_similarity", 0.2)
        search_type = self.config.get("search_type", "hybrid")
        
        # Default empty results
        results = []
        
        if self.collection:
            # Convert filter criteria to Chroma format if needed
            where_clause = {}
            if filter_criteria:
                where_clause = filter_criteria
            
            # Get the main query embedding
            query_embedding = self.embedder.get_embedding(query)
            
            if search_type == "semantic" or search_type == "hybrid":
                # Perform vector search for each processed query
                all_semantic_results = []
                
                for processed_query in processed_queries:
                    # Get embedding for this query
                    query_embedding = self.embedder.get_embedding(processed_query)
                    
                    # Skip this query if embedding failed
                    if query_embedding is None:
                        continue
                    
                    # Convert numpy array to list for ChromaDB
                    query_embedding_list = query_embedding.tolist() if isinstance(query_embedding, np.ndarray) else query_embedding
                    
                    # Get results from vector database
                    semantic_results = self.collection.query(
                        query_embeddings=[query_embedding_list],
                        n_results=top_k * 2,  # Get more results than needed for re-ranking
                        where=where_clause
                    )
                    
                    # Process results
                    for i in range(len(semantic_results['ids'][0])):
                        result = {
                            "chunk_id": semantic_results['ids'][0][i],
                            "content": semantic_results['documents'][0][i],
                            "metadata": semantic_results['metadatas'][0][i],
                            "similarity": semantic_results['distances'][0][i] if 'distances' in semantic_results else 0.0,
                            "source": "semantic"
                        }
                        all_semantic_results.append(result)
                
                # Add semantic results to final results
                results.extend(all_semantic_results)
            
            if search_type == "keyword" or search_type == "hybrid":
                # Perform keyword search
                keyword_results = []
                
                for processed_query in processed_queries:
                    try:
                        # Use Chroma's built-in keyword search if available
                        keyword_response = self.collection.query(
                            query_texts=[processed_query],
                            n_results=top_k * 2,
                            where=where_clause
                        )
                        
                        # Process results
                        for i in range(len(keyword_response['ids'][0])):
                            result = {
                                "chunk_id": keyword_response['ids'][0][i],
                                "content": keyword_response['documents'][0][i],
                                "metadata": keyword_response['metadatas'][0][i],
                                "similarity": keyword_response['distances'][0][i] if 'distances' in keyword_response else 0.0,
                                "source": "keyword"
                            }
                            keyword_results.append(result)
                    except Exception as e:
                        logger.error(f"Error performing keyword search: {str(e)}")
                
                # Add keyword results to final results
                results.extend(keyword_results)
            
            # Remove duplicates
            unique_results = {}
            for result in results:
                chunk_id = result["chunk_id"]
                if chunk_id not in unique_results or result["similarity"] > unique_results[chunk_id]["similarity"]:
                    unique_results[chunk_id] = result
            
            results = list(unique_results.values())
            
            # Apply minimum similarity threshold
            results = [result for result in results if result["similarity"] >= min_similarity]
            
            # Re-rank results if configured
            if self.config.get("use_reranking", True):
                results = self._rerank_results(query, results)
            
            # Sort by similarity
            results.sort(key=lambda x: x["similarity"], reverse=True)
            
            # Limit to top_k
            results = results[:top_k]
            
            # Group by document and maintain order
            grouped_results = defaultdict(list)
            for result in results:
                doc_id = result["metadata"]["doc_id"]
                grouped_results[doc_id].append(result)
            
            # Get document-level metadata and add to results
            final_results = []
            for doc_id, chunks in grouped_results.items():
                doc_metadata = self.document_metadata.get(doc_id, {})
                
                # Add document metadata to first chunk
                for chunk in chunks:
                    chunk["document_metadata"] = {
                        "doc_id": doc_id,
                        "title": doc_metadata.get("title", ""),
                        "summary": doc_metadata.get("summary", ""),
                        "category": doc_metadata.get("category", ""),
                        "source": doc_metadata.get("source", ""),
                        "doc_type": doc_metadata.get("doc_type", "")
                    }
                    final_results.append(chunk)
            
            return final_results
        else:
            # If not using a vector database, we can't perform search
            logger.warning("Cannot perform search without vector database")
            return []
    
    def _rerank_results(self, query: str, results: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Re-rank search results using a more sophisticated approach.
        
        In a production environment, this would use a dedicated re-ranker model.
        For now, we'll use a simple heuristic.
        """
        reranker = self.config.get("reranker", "default")
        
        if reranker == "default":
            # Simple re-ranking: boost results that contain query terms in title
            query_terms = query.lower().split()
            
            for result in results:
                title = result.get("metadata", {}).get("document_title", "").lower()
                boost = 0
                
                # Check if query terms appear in title
                for term in query_terms:
                    if term in title:
                        boost += 0.1
                
                # Check if query terms appear in content
                content = result.get("content", "").lower()
                for term in query_terms:
                    if term in content:
                        boost += 0.05
                
                # Apply boost
                result["similarity"] = min(1.0, result["similarity"] + boost)
        
        return results


class AdvancedRAG:
    """Main class for advanced RAG operations."""
    
    def __init__(self, config_name: str = "default"):
        """Initialize the advanced RAG system."""
        self.config_name = config_name
        self.config = self._load_config()
        
        self.document_processor = DocumentProcessor({
            "chunking": self.config["chunking"], 
            "document_processing": self.config["document_processing"]
        })
        
        self.document_store = DocumentStore({
            "vector_db": self.config["vector_db"],
            "retrieval": self.config["retrieval"],
            "embedding": self.config["embedding"],
            "query_processing": self.config["query_processing"]
        })
    
    def _load_config(self) -> Dict[str, Any]:
        """Load configuration from file or use default."""
        config_file = os.path.join(CONFIG_DIR, f"{self.config_name}.json")
        
        if os.path.exists(config_file):
            try:
                with open(config_file, 'r') as f:
                    config = json.load(f)
                    
                # Merge with default config to ensure all fields are present
                merged_config = DEFAULT_CONFIG.copy()
                for key, value in config.items():
                    if key in merged_config and isinstance(merged_config[key], dict) and isinstance(value, dict):
                        merged_config[key].update(value)
                    else:
                        merged_config[key] = value
                        
                return merged_config
            except Exception as e:
                logger.error(f"Error loading config: {str(e)}")
                return DEFAULT_CONFIG
        else:
            # Create default config
            try:
                os.makedirs(os.path.dirname(config_file), exist_ok=True)
                with open(config_file, 'w') as f:
                    json.dump(DEFAULT_CONFIG, f, indent=2)
            except Exception as e:
                logger.error(f"Error creating default config: {str(e)}")
                
            return DEFAULT_CONFIG
    
    def save_config(self) -> None:
        """Save current configuration to file."""
        config_file = os.path.join(CONFIG_DIR, f"{self.config_name}.json")
        
        try:
            with open(config_file, 'w') as f:
                json.dump(self.config, f, indent=2)
                
            logger.info(f"Configuration saved to {config_file}")
        except Exception as e:
            logger.error(f"Error saving configuration: {str(e)}")
    
    def add_document(self, file_path: str) -> str:
        """Add a document from a file to the RAG system."""
        # Process document
        document = self.document_processor.process_file(file_path)
        
        # Add to document store
        doc_id = self.document_store.add_document(document)
        
        return doc_id
    
    def add_text(self, text: str, title: str = None, source: str = None) -> str:
        """Add text directly to the RAG system."""
        # Process text
        document = self.document_processor.process_text(text, title, source)
        
        # Add to document store
        doc_id = self.document_store.add_document(document)
        
        return doc_id
    
    def search(self, query: str, top_k: int = None, filter_criteria: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """Search for documents matching the query."""
        return self.document_store.search(query, top_k, filter_criteria=filter_criteria)
    
    def retrieve_context(self, query: str, top_k: int = 3, filter_criteria: Dict[str, Any] = None) -> str:
        """Retrieve context for a query in a format suitable for RAG."""
        results = self.search(query, top_k, filter_criteria=filter_criteria)
        
        if not results:
            return ""
        
        context = []
        for i, result in enumerate(results):
            doc_title = result.get("document_metadata", {}).get("title", "Untitled Document")
            doc_source = result.get("document_metadata", {}).get("source", "Unknown Source")
            
            context.append(f"[Document {i+1}: {doc_title} (Source: {doc_source})]")
            context.append(result["content"])
            context.append("")
        
        return "\n".join(context)
    
    def get_document(self, doc_id: str) -> Optional[Document]:
        """Get a document by ID."""
        return self.document_store.get_document(doc_id)
    
    def delete_document(self, doc_id: str) -> bool:
        """Delete a document by ID."""
        return self.document_store.delete_document(doc_id)
    
    def list_documents(self) -> List[Dict[str, Any]]:
        """List all documents in the store."""
        return self.document_store.list_documents()


def rag_ui():
    """Streamlit UI for the Advanced RAG system."""
    st.title("Advanced RAG")
    st.subheader("Retrieval Augmented Generation")
    
    # Create tabs
    tabs = st.tabs(["Search", "Documents", "Upload", "Settings"])
    
    # Initialize RAG system
    rag = AdvancedRAG()
    
    with tabs[0]:  # Search tab
        st.header("Search Documents")
        
        # Search interface
        query = st.text_input("Enter your query:", key="search_query")
        
        # Filter options
        with st.expander("Search Filters", expanded=False):
            # Category filter
            documents = rag.list_documents()
            categories = sorted(list(set([doc.get("category", "Uncategorized") for doc in documents])))
            
            selected_categories = st.multiselect(
                "Filter by category:",
                options=categories,
                default=None,
                key="search_categories"
            )
            
            # Document type filter
            doc_types = sorted(list(set([doc.get("doc_type", "Unknown") for doc in documents])))
            
            selected_doc_types = st.multiselect(
                "Filter by document type:",
                options=doc_types,
                default=None,
                key="search_doc_types"
            )
            
            # Source filter
            sources = sorted(list(set([doc.get("source", "Unknown") for doc in documents])))
            
            selected_sources = st.multiselect(
                "Filter by source:",
                options=sources,
                default=None,
                key="search_sources"
            )
            
            # Number of results
            top_k = st.slider("Number of results:", min_value=1, max_value=20, value=5, step=1, key="search_top_k")
        
        # Build filter criteria
        filter_criteria = {}
        if selected_categories:
            filter_criteria["category"] = {"$in": selected_categories}
        if selected_doc_types:
            filter_criteria["doc_type"] = {"$in": selected_doc_types}
        if selected_sources:
            filter_criteria["source"] = {"$in": selected_sources}
        
        # Search button
        if st.button("Search", key="search_button") and query:
            with st.spinner("Searching..."):
                results = rag.search(query, top_k=top_k, filter_criteria=filter_criteria if filter_criteria else None)
                
                if results:
                    st.success(f"Found {len(results)} results")
                    
                    # Display results
                    for i, result in enumerate(results):
                        with st.container(border=True):
                            # Get document metadata
                            doc_metadata = result.get("document_metadata", {})
                            doc_title = doc_metadata.get("title", "Untitled Document")
                            doc_source = doc_metadata.get("source", "Unknown Source")
                            doc_type = doc_metadata.get("doc_type", "Unknown Type")
                            doc_category = doc_metadata.get("category", "Uncategorized")
                            
                            # Header with title and metadata
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                st.subheader(doc_title)
                            with col2:
                                st.caption(f"Score: {result.get('similarity', 0):.4f}")
                            
                            # Content
                            st.markdown(result["content"])
                            
                            # Footer with metadata
                            st.caption(f"Source: {doc_source} | Type: {doc_type} | Category: {doc_category}")
                            
                            # Option to view full document
                            doc_id = doc_metadata.get("doc_id")
                            if doc_id and st.button(f"View Full Document", key=f"view_doc_{i}_{doc_id}"):
                                st.session_state.selected_doc_id = doc_id
                                st.switch_page("Advanced RAG", "Documents")
                else:
                    st.warning("No results found")
        
        # RAG context generation
        with st.expander("Generate RAG Context", expanded=False):
            context_query = st.text_input("Enter query for context generation:", value=query, key="context_query")
            context_limit = st.slider("Number of documents:", min_value=1, max_value=10, value=3, step=1, key="context_limit")
            
            if st.button("Generate Context", key="generate_context_button") and context_query:
                with st.spinner("Generating context..."):
                    context = rag.retrieve_context(context_query, top_k=context_limit, filter_criteria=filter_criteria)
                    
                    if context:
                        st.subheader("Generated Context")
                        st.text_area("Context (copy this to use with LLMs):", value=context, height=300)
                        
                        # Option to copy to clipboard (JavaScript)
                        st.markdown("""
                        <button onclick="navigator.clipboard.writeText(document.querySelector('textarea').value)">
                            Copy to Clipboard
                        </button>
                        """, unsafe_allow_html=True)
                    else:
                        st.warning("No context generated")
    
    with tabs[1]:  # Documents tab
        st.header("Manage Documents")
        
        # List documents
        documents = rag.list_documents()
        
        if documents:
            # Create a dataframe for better display
            docs_data = []
            for doc in documents:
                docs_data.append({
                    "ID": doc.get("doc_id", ""),
                    "Title": doc.get("title", "Untitled"),
                    "Type": doc.get("doc_type", "Unknown"),
                    "Category": doc.get("category", "Uncategorized"),
                    "Chunks": doc.get("chunk_count", 0),
                    "Created": doc.get("created_at", "")
                })
            
            # Convert to dataframe
            df = pd.DataFrame(docs_data)
            
            # Display with filters
            with st.container(height=400):
                st.dataframe(df, use_container_width=True)
            
            # Document selection
            selected_doc_id = st.selectbox(
                "Select a document to view or delete:",
                options=[doc.get("doc_id", "") for doc in documents],
                format_func=lambda x: next((doc.get("title", "Untitled") for doc in documents if doc.get("doc_id", "") == x), x),
                key="doc_selector"
            )
            
            if selected_doc_id:
                # Get document
                document = rag.get_document(selected_doc_id)
                
                if document:
                    # Document actions
                    col1, col2 = st.columns([1, 1])
                    with col1:
                        if st.button("View Document", key="view_document_button"):
                            st.subheader(document.title)
                            st.caption(f"Source: {document.source} | Type: {document.doc_type} | Category: {document.category}")
                            
                            with st.expander("Document Metadata", expanded=False):
                                st.json(document.metadata)
                            
                            st.markdown(document.content)
                    with col2:
                        if st.button("Delete Document", key="delete_document_button"):
                            if rag.delete_document(selected_doc_id):
                                st.success(f"Document '{document.title}' deleted successfully")
                                time.sleep(1)
                                st.rerun()
                            else:
                                st.error("Failed to delete document")
                else:
                    st.error(f"Document {selected_doc_id} not found")
        else:
            st.info("No documents found. Upload documents in the Upload tab.")
    
    with tabs[2]:  # Upload tab
        st.header("Upload Documents")
        
        # File uploader
        uploaded_file = st.file_uploader("Upload a document:", type=["txt", "pdf", "html", "md", "docx", "json", "csv"])
        
        if uploaded_file is not None:
            # Save the uploaded file
            file_path = os.path.join(DOCUMENT_STORE_DIR, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            st.success(f"File uploaded: {uploaded_file.name}")
            
            # Process the file
            if st.button("Process Document", key="process_document_button"):
                with st.spinner("Processing document..."):
                    try:
                        doc_id = rag.add_document(file_path)
                        st.success(f"Document processed and added with ID: {doc_id}")
                    except Exception as e:
                        st.error(f"Error processing document: {str(e)}")
        
        # Text input
        st.subheader("Or add text directly:")
        
        text_title = st.text_input("Title:", key="text_title")
        text_source = st.text_input("Source (optional):", key="text_source")
        text_content = st.text_area("Content:", height=200, key="text_content")
        
        if st.button("Add Text", key="add_text_button") and text_content:
            with st.spinner("Processing text..."):
                try:
                    doc_id = rag.add_text(text_content, text_title, text_source)
                    st.success(f"Text added with ID: {doc_id}")
                except Exception as e:
                    st.error(f"Error adding text: {str(e)}")
    
    with tabs[3]:  # Settings tab
        st.header("RAG Settings")
        
        # Display current configuration
        with st.expander("Current Configuration", expanded=True):
            st.json(rag.config)
        
        # Allow modifying configuration
        st.subheader("Modify Configuration")
        
        # Chunking settings
        st.markdown("#### Chunking Settings")
        
        chunking_method = st.selectbox(
            "Chunking Method:",
            options=["semantic", "fixed", "sliding", "paragraph"],
            index=["semantic", "fixed", "sliding", "paragraph"].index(rag.config["chunking"]["method"]),
            key="chunking_method"
        )
        
        chunk_size = st.slider(
            "Chunk Size:",
            min_value=100,
            max_value=5000,
            value=rag.config["chunking"]["size"],
            step=100,
            key="chunk_size"
        )
        
        chunk_overlap = st.slider(
            "Chunk Overlap:",
            min_value=0,
            max_value=min(chunk_size, 1000),
            value=min(rag.config["chunking"]["overlap"], min(chunk_size, 1000)),
            step=50,
            key="chunk_overlap"
        )
        
        # Embedding settings
        st.markdown("#### Embedding Settings")
        
        embedding_provider = st.selectbox(
            "Embedding Provider:",
            options=["ollama", "openai"],
            index=["ollama", "openai"].index(rag.config["embedding"]["provider"]) if rag.config["embedding"]["provider"] in ["ollama", "openai"] else 0,
            key="embedding_provider"
        )
        
        embedding_model = st.text_input(
            "Embedding Model:",
            value=rag.config["embedding"]["model"],
            key="embedding_model"
        )
        
        cache_embeddings = st.checkbox(
            "Cache Embeddings",
            value=rag.config["embedding"]["cache_embeddings"],
            key="cache_embeddings"
        )
        
        # Retrieval settings
        st.markdown("#### Retrieval Settings")
        
        search_type = st.selectbox(
            "Search Type:",
            options=["semantic", "keyword", "hybrid"],
            index=["semantic", "keyword", "hybrid"].index(rag.config["retrieval"]["search_type"]) if rag.config["retrieval"]["search_type"] in ["semantic", "keyword", "hybrid"] else 0,
            key="search_type"
        )
        
        if search_type == "hybrid":
            col1, col2 = st.columns(2)
            with col1:
                semantic_weight = st.slider(
                    "Semantic Weight:",
                    min_value=0.0,
                    max_value=1.0,
                    value=rag.config["retrieval"]["semantic_weight"],
                    step=0.1,
                    key="semantic_weight"
                )
            with col2:
                keyword_weight = st.slider(
                    "Keyword Weight:",
                    min_value=0.0,
                    max_value=1.0,
                    value=rag.config["retrieval"]["keyword_weight"],
                    step=0.1,
                    key="keyword_weight"
                )
        
        default_top_k = st.slider(
            "Default Number of Results:",
            min_value=1,
            max_value=20,
            value=rag.config["retrieval"]["top_k"],
            step=1,
            key="default_top_k"
        )
        
        min_similarity = st.slider(
            "Minimum Similarity Score:",
            min_value=0.0,
            max_value=1.0,
            value=rag.config["retrieval"]["min_similarity"],
            step=0.05,
            key="min_similarity"
        )
        
        use_reranking = st.checkbox(
            "Use Reranking",
            value=rag.config["retrieval"]["use_reranking"],
            key="use_reranking"
        )
        
        # Save settings
        if st.button("Save Settings", key="save_settings_button"):
            # Update configuration
            rag.config["chunking"]["method"] = chunking_method
            rag.config["chunking"]["size"] = chunk_size
            rag.config["chunking"]["overlap"] = chunk_overlap
            
            rag.config["embedding"]["provider"] = embedding_provider
            rag.config["embedding"]["model"] = embedding_model
            rag.config["embedding"]["cache_embeddings"] = cache_embeddings
            
            rag.config["retrieval"]["search_type"] = search_type
            if search_type == "hybrid":
                rag.config["retrieval"]["semantic_weight"] = semantic_weight
                rag.config["retrieval"]["keyword_weight"] = keyword_weight
            rag.config["retrieval"]["top_k"] = default_top_k
            rag.config["retrieval"]["min_similarity"] = min_similarity
            rag.config["retrieval"]["use_reranking"] = use_reranking
            
            # Save configuration
            rag.save_config()
            
            st.success("Settings saved successfully")
            time.sleep(1)
            st.rerun()


if __name__ == "__main__":
    rag_ui()